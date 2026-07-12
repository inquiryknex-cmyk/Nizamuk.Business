#!/usr/bin/env python3
"""
Convert a designed PDF into a visually-faithful, fully text-editable PPTX.

Strategy per page:
  1. Remove ONLY the text layer (keep images, logos, vector art) and render
     the rest as a high-res background image for the slide.
  2. Re-create every text line as native, editable PowerPoint text boxes at
     the exact same position, size and color.
Pages with a size different from the dominant page size are scaled to fit
and centered (slide background filled with the page edge color).
"""
import io, os, math, sys
import fitz
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

SRC = "/root/.claude/uploads/4816b3b7-f4dc-583c-b56e-ed0b4a3e381d/033c66b4-FM__FS_Summit_2026_Riyadh__KE.pdf"
OUT = sys.argv[1] if len(sys.argv) > 1 else "FM_FS_Summit_2026_Riyadh_KE.pptx"
WORK = os.path.dirname(os.path.abspath(OUT)) or "."
EMU_PER_PT = 12700
TARGET_LONG_PX = 2800          # background render resolution (long side)
GAP_FACTOR = 1.9               # split a line into segments when the gap
                               # between spans exceeds GAP_FACTOR * space width

def map_font(pdf_font, flags):
    name = pdf_font.split("+")[-1]
    low = name.lower()
    bold = bool(flags & 16) or any(k in low for k in ("bold", "semibold", "demi", "heavy", "black", "bolder"))
    italic = bool(flags & 2) or "italic" in low or "oblique" in low
    if "wingding" in low:                       fam = "Wingdings"
    elif low.startswith("mont"):                fam = "Montserrat"
    elif "avantgarde" in low:                   fam = "Century Gothic"
    elif "futura" in low:                       fam = "Century Gothic"
    elif "gillsans" in low:                     fam = "Gill Sans MT"
    elif "arial" in low:                        fam = "Arial"
    elif "janna" in low:                        fam = "Janna LT"
    elif "ebrima" in low:                       fam = "Ebrima"
    elif "extended" in low:                     fam = "Agency FB"      # condensed display headers
    elif "infinitestroke" in low:               fam = "Cooper Black"   # heavy display banner
    else:                                       fam = name
    return fam, bold, italic

def int_to_rgb(c):
    return RGBColor((c >> 16) & 255, (c >> 8) & 255, c & 255)

def collect_segments(page):
    """Return list of segments: dicts with bbox, dir, list of runs."""
    d = page.get_text("dict", flags=fitz.TEXTFLAGS_DICT & ~fitz.TEXT_PRESERVE_LIGATURES)
    segments = []
    for b in d["blocks"]:
        if b["type"] != 0:
            continue
        for l in b["lines"]:
            spans = [s for s in l["spans"] if s["text"].strip("  \t")]
            if not spans:
                continue
            direction = l["dir"]
            cur = None
            for s in spans:
                sb = fitz.Rect(s["bbox"])
                if cur is not None:
                    gap_lim = GAP_FACTOR * max(2.0, cur["last_size"] * 0.5)
                    horizontal = abs(direction[0]) > 0.7
                    gap = (sb.x0 - cur["bbox"].x1) if horizontal else (sb.y0 - cur["bbox"].y1)
                    if gap > gap_lim:
                        segments.append(cur); cur = None
                if cur is None:
                    cur = {"bbox": fitz.Rect(sb), "dir": direction, "runs": [], "last_size": s["size"]}
                else:
                    cur["bbox"] |= sb
                    cur["last_size"] = s["size"]
                cur["runs"].append(s)
            if cur:
                segments.append(cur)
    # de-duplicate drop-shadow twins (same text drawn twice, tiny offset) — keep last (on top)
    keep = [True] * len(segments)
    def seg_text(seg): return "".join(r["text"] for r in seg["runs"]).strip()
    for i in range(len(segments)):
        for j in range(i + 1, len(segments)):
            if not keep[i]:
                break
            a, bb = segments[i], segments[j]
            if seg_text(a) == seg_text(bb):
                ca, cb = (a["bbox"].tl + a["bbox"].br) / 2, (bb["bbox"].tl + bb["bbox"].br) / 2
                tol = 0.45 * max(a["last_size"], 2)
                if abs(ca.x - cb.x) < tol and abs(ca.y - cb.y) < tol:
                    keep[i] = False
    return [s for k, s in zip(keep, segments) if k]

def strip_text_and_render(doc_path, pno, zoom):
    doc = fitz.open(doc_path)
    page = doc[pno]
    for b in page.get_text("dict")["blocks"]:
        if b["type"] != 0:
            continue
        for l in b["lines"]:
            for s in l["spans"]:
                r = fitz.Rect(s["bbox"]) & page.rect
                if not r.is_empty:
                    page.add_redact_annot(r)
    try:
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE,
                              graphics=fitz.PDF_REDACT_LINE_ART_NONE,
                              text=fitz.PDF_REDACT_TEXT_REMOVE)
    except TypeError:
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
    pm = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    img = Image.frombytes("RGB", (pm.width, pm.height), pm.samples)
    doc.close()
    return img

def edge_color(img):
    px = img.load()
    w, h = img.size
    pts = [(3, 3), (w - 4, 3), (3, h - 4), (w - 4, h - 4)]
    r = sum(px[p][0] for p in pts) // 4
    g = sum(px[p][1] for p in pts) // 4
    b = sum(px[p][2] for p in pts) // 4
    return RGBColor(r, g, b)

def main():
    doc = fitz.open(SRC)
    # dominant page size -> slide size
    from collections import Counter
    cnt = Counter((round(p.rect.width, 1), round(p.rect.height, 1)) for p in doc)
    (SW, SH), _ = cnt.most_common(1)[0]
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_PER_PT))
    prs.slide_height = Emu(int(SH * EMU_PER_PT))
    blank = prs.slide_layouts[6]

    for pno in range(len(doc)):
        page = doc[pno]
        pw, ph = page.rect.width, page.rect.height
        scale = min(SW / pw, SH / ph)
        ox, oy = (SW - pw * scale) / 2.0, (SH - ph * scale) / 2.0

        slide = prs.slides.add_slide(blank)
        zoom = TARGET_LONG_PX / max(pw, ph)
        img = strip_text_and_render(SRC, pno, zoom)

        if scale < 0.999:  # letterboxed page: paint slide bg with page edge color
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = edge_color(img)

        # save background (JPEG q92; PNG if smaller)
        jp, pn = io.BytesIO(), io.BytesIO()
        img.save(jp, "JPEG", quality=92, optimize=True)
        img.save(pn, "PNG", optimize=True)
        buf, ext = (pn, "png") if pn.tell() < jp.tell() else (jp, "jpg")
        ipath = os.path.join(WORK, f"bg_{pno+1:02d}.{ext}")
        with open(ipath, "wb") as f:
            f.write(buf.getvalue())
        slide.shapes.add_picture(ipath,
                                 Emu(int(ox * EMU_PER_PT)), Emu(int(oy * EMU_PER_PT)),
                                 Emu(int(pw * scale * EMU_PER_PT)), Emu(int(ph * scale * EMU_PER_PT)))

        for seg in collect_segments(page):
            bb = seg["bbox"]
            vertical = abs(seg["dir"][1]) > 0.7
            x0, y0 = bb.x0 * scale + ox, bb.y0 * scale + oy
            w, h = bb.width * scale, bb.height * scale
            if vertical:
                cx, cy = x0 + w / 2, y0 + h / 2
                w, h = h, w
                x0, y0 = cx - w / 2, cy - h / 2
            # a little slack so substituted fonts don't clip
            pad_w = w * 0.30 + 4
            tb = slide.shapes.add_textbox(Emu(int(x0 * EMU_PER_PT)), Emu(int(y0 * EMU_PER_PT)),
                                          Emu(int((w + pad_w) * EMU_PER_PT)), Emu(int(max(h, 2) * EMU_PER_PT)))
            if vertical:
                # rotate around center; text reads bottom-to-top
                tb.rotation = 270 if seg["dir"][1] < 0 else 90
                tb.left = Emu(int((x0 - pad_w / 2) * EMU_PER_PT))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            has_arabic = any("؀" <= ch <= "ۿ" for r in seg["runs"] for ch in r["text"])
            if has_arabic:
                p._p.get_or_add_pPr().set("rtl", "1")
                p.alignment = PP_ALIGN.RIGHT
            for s in seg["runs"]:
                run = p.add_run()
                run.text = s["text"].rstrip("\n")
                fam, bold, italic = map_font(s["font"], s["flags"])
                fs = max(1.0, s["size"] * scale)
                run.font.size = Pt(round(fs * 2) / 2)
                run.font.name = fam
                run.font.bold = bold
                run.font.italic = italic
                run.font.color.rgb = int_to_rgb(s["color"])
                if s["flags"] & 1:  # superscript
                    run.font._rPr.set("baseline", "30000")

    prs.save(OUT)
    print("saved", OUT, f"{os.path.getsize(OUT)/1e6:.1f} MB, {len(doc)} slides")

if __name__ == "__main__":
    main()
