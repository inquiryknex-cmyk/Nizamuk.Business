#!/usr/bin/env python3
"""
Convert a designed portrait PDF into a 16:9 landscape PPTX that FILLS each
slide: every page is auto-segmented into content blocks at natural section
boundaries, blocks are re-packed into balanced columns across the widescreen
canvas, graphics stay pixel-perfect (cropped from a text-stripped render)
and all text is re-created as editable text boxes at the new positions.
"""
import io, os, sys
import fitz
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

SRC = "/root/.claude/uploads/4816b3b7-f4dc-583c-b56e-ed0b4a3e381d/033c66b4-FM__FS_Summit_2026_Riyadh__KE.pdf"
OUT = sys.argv[1] if len(sys.argv) > 1 else "FM_FS_Summit_2026_Riyadh_KE_wide.pptx"
WORK = os.path.dirname(os.path.abspath(OUT)) or "."
EMU = 12700
SW, SH = 960.0, 540.0            # 16:9 slide in points
SKIP_PAGES = {10, 11}
MARGIN, GUTTER = 14.0, 16.0
RENDER_LONG_PX = 3200
Y_GAP = 7.0                      # min vertical whitespace to cut a band
X_GAP = 12.0                     # min horizontal whitespace for grid split
PAD = 4.0                        # padding kept around each cropped block
S_CAP = 1.25                     # never upscale content beyond this

DIST_THRESH = 55               # pixel distance from bg color to count as content

def map_font(pdf_font, flags):
    name = pdf_font.split("+")[-1]
    low = name.lower()
    bold = bool(flags & 16) or any(k in low for k in ("bold", "semibold", "demi", "heavy", "black", "bolder"))
    italic = bool(flags & 2) or "italic" in low or "oblique" in low
    if "wingding" in low:         fam = "Wingdings"
    elif low.startswith("mont"):  fam = "Montserrat"
    elif "avantgarde" in low:     fam = "Century Gothic"
    elif "futura" in low:         fam = "Century Gothic"
    elif "gillsans" in low:       fam = "Gill Sans MT"
    elif "arial" in low:          fam = "Arial"
    elif "janna" in low:          fam = "Janna LT"
    elif "ebrima" in low:         fam = "Ebrima"
    elif "extended" in low:       fam = "Agency FB"
    elif "infinitestroke" in low: fam = "Cooper Black"
    else:                         fam = name
    return fam, bold, italic

def int_to_rgb(c):
    return RGBColor((c >> 16) & 255, (c >> 8) & 255, c & 255)

def luminance(color):
    if color is None:
        return 1.0
    if isinstance(color, (tuple, list)):
        r, g, b = color[:3]
    else:
        r, g, b = ((color >> 16) & 255) / 255, ((color >> 8) & 255) / 255, (color & 255) / 255
        return 0.299 * r + 0.587 * g + 0.114 * b
    return 0.299 * r + 0.587 * g + 0.114 * b

def collect_segments(page):
    d = page.get_text("dict", flags=fitz.TEXTFLAGS_DICT & ~fitz.TEXT_PRESERVE_LIGATURES)
    segments = []
    for b in d["blocks"]:
        if b["type"] != 0:
            continue
        for l in b["lines"]:
            spans = [s for s in l["spans"] if s["text"].strip("  \t")]
            if not spans:
                continue
            direction = l["dir"]
            cur = None
            for s in spans:
                sb = fitz.Rect(s["bbox"])
                if cur is not None:
                    gap_lim = 1.9 * max(2.0, cur["last_size"] * 0.5)
                    horizontal = abs(direction[0]) > 0.7
                    gap = (sb.x0 - cur["bbox"].x1) if horizontal else (sb.y0 - cur["bbox"].y1)
                    if gap > gap_lim:
                        segments.append(cur); cur = None
                if cur is None:
                    cur = {"bbox": fitz.Rect(sb), "dir": direction, "runs": [], "last_size": s["size"]}
                else:
                    cur["bbox"] |= sb
                    cur["last_size"] = s["size"]
                cur["runs"].append(s)
            if cur:
                segments.append(cur)
    keep = [True] * len(segments)
    def seg_text(seg): return "".join(r["text"] for r in seg["runs"]).strip()
    for i in range(len(segments)):
        for j in range(i + 1, len(segments)):
            if not keep[i]:
                break
            a, bb = segments[i], segments[j]
            if seg_text(a) == seg_text(bb):
                ca, cb = (a["bbox"].tl + a["bbox"].br) / 2, (bb["bbox"].tl + bb["bbox"].br) / 2
                tol = 0.45 * max(a["last_size"], 2)
                if abs(ca.x - cb.x) < tol and abs(ca.y - cb.y) < tol:
                    keep[i] = False
    return [s for k, s in zip(keep, segments) if k]

import numpy as np

# per-page tuning overrides (0-based page number)
PAGE_OPTS = {
    2: {"people_grid": True},
    5: {"y_gap": 3},
    6: {"forced": ["Speaking Opportunities", "Networking & Delegate Access",
                   "Exhibit Space & Onsite Presence",
                   "Post-Event Recognition & Access", "Special Honors"],
        "y_gap": 14},
    8: {"dist": 75, "edge_zero": 50, "cut_tol": 14},
}

def occupancy(page, img_hi, zoom, dist_thresh, edge_zero=0):
    """Boolean occupancy grid at ~1px per PDF point, from the text-stripped
    render (pale background texture ignored) plus stamped text bboxes."""
    pw, ph = int(round(page.rect.width)), int(round(page.rect.height))
    small = img_hi.resize((pw, ph), Image.BILINEAR)
    arr = np.asarray(small, dtype=np.int16)
    ez = edge_zero
    interior = arr[ez:ph - ez, ez:pw - ez] if ez else arr
    q = (interior // 16) * 16 + 8
    colors, counts = np.unique(q.reshape(-1, 3), axis=0, return_counts=True)
    bg = colors[counts.argmax()]
    dist = np.abs(arr - bg).max(axis=2)
    occ = dist > dist_thresh
    if edge_zero:
        occ[:edge_zero] = occ[-edge_zero:] = False
        occ[:, :edge_zero] = occ[:, -edge_zero:] = False
    for seg_b in page.get_text("dict")["blocks"]:
        if seg_b["type"] != 0:
            continue
        for l in seg_b["lines"]:
            for s in l["spans"]:
                if not s["text"].strip():
                    continue
                x0, y0, x1, y1 = s["bbox"]
                occ[max(0, int(y0)):min(ph, int(y1) + 1),
                    max(0, int(x0)):min(pw, int(x1) + 1)] = True
    return occ

def dilate(mask, n, axis):
    out = mask.copy()
    for k in range(1, n + 1):
        out[:-k] |= mask[k:] if axis == 0 else False
        out[k:] |= mask[:-k] if axis == 0 else False
        if axis == 1:
            out[:, :-k] |= mask[:, k:]
            out[:, k:] |= mask[:, :-k]
    return out

def runs_of(cond):
    """[(start, end)] of maximal True runs."""
    idx = np.flatnonzero(np.diff(np.concatenate(([False], cond, [False])).astype(np.int8)))
    return list(zip(idx[::2], idx[1::2]))

def find_cuts(page, occ, y_gap, forced_texts, cut_tol=2):
    """Return y cut positions. A cut needs a horizontal corridor that is
    either truly empty or only crossed by thin full-width separator lines;
    vertical box borders (narrow occupied columns) block cuts."""
    vdil = dilate(occ, 2, axis=0)
    row_cnt = vdil.sum(axis=1)
    width_occ = occ.any(axis=0).sum()
    cuttable = row_cnt <= cut_tol
    # thin, wide separator lines: short non-cuttable runs with high coverage
    separator = np.zeros_like(cuttable)
    for a, b in runs_of(~cuttable):
        if b - a <= 6 and row_cnt[a:b].max() >= 0.20 * width_occ:
            separator[a:b] = True
    corridor = cuttable | separator
    cuts = []
    for a, b in runs_of(corridor):
        if b - a < y_gap:
            continue
        sub = [(sa, sb) for sa, sb in runs_of(cuttable[a:b]) if sb - sa >= 2]
        if not sub:
            continue
        sa, sb = max(sub, key=lambda r: r[1] - r[0])
        cuts.append(a + (sa + sb) // 2)
    # forced cuts at named section headers
    for b in page.get_text("dict")["blocks"]:
        if b["type"] != 0:
            continue
        for l in b["lines"]:
            line_text = "".join(s["text"] for s in l["spans"]).strip()
            if any(line_text.startswith(t) for t in forced_texts):
                cuts.append(int(l["bbox"][1] - 2))
    return sorted(set(cuts))

def trim_rect(occ, rect, pad=PAD):
    """Tighten rect to occupied content inside it, then pad."""
    ph, pw = occ.shape
    y0, y1 = max(0, int(rect.y0)), min(ph, int(rect.y1))
    x0, x1 = max(0, int(rect.x0)), min(pw, int(rect.x1))
    sub = occ[y0:y1, x0:x1]
    if not sub.any():
        return None
    rows, cols = sub.any(axis=1), sub.any(axis=0)
    ty0 = y0 + int(np.argmax(rows)); ty1 = y1 - int(np.argmax(rows[::-1]))
    tx0 = x0 + int(np.argmax(cols)); tx1 = x1 - int(np.argmax(cols[::-1]))
    return fitz.Rect(max(0, tx0 - pad), max(0, ty0 - pad),
                     min(pw, tx1 + pad), min(ph, ty1 + pad))

def people_grid_blocks(page, occ):
    """Executive-board style page: one block per person (photo + name +
    titles), plus separately-scaled banner blocks for the headings."""
    d = page.get_text("dict")
    ext_lines, other_rects = [], []
    for b in d["blocks"]:
        if b["type"] != 0:
            continue
        for l in b["lines"]:
            spans = [s for s in l["spans"] if s["text"].strip()]
            if not spans:
                continue
            r = fitz.Rect(l["bbox"])
            text = "".join(s["text"] for s in spans).strip()
            if any("Extended" in s["font"] for s in spans):
                ext_lines.append((r, text))
            else:
                other_rects.append(r)
    names = []
    for r, t in ext_lines:
        if r.y0 < 110:
            continue
        # a person name has a plain-font line (title/role) right below it
        if any(-5 < o.y0 - r.y1 < 25 and o.x0 < r.x1 and o.x1 > r.x0
               for o in other_rects):
            names.append(r)
    if len(names) < 4:
        return None
    names.sort(key=lambda r: (r.y0, r.x0))
    rows, cur = [], [names[0]]
    for r in names[1:]:
        if abs(r.y0 - cur[0].y0) < 40:
            cur.append(r)
        else:
            rows.append(cur); cur = [r]
    rows.append(cur)
    ph, pw = occ.shape
    # vertical boundaries between people rows
    bounds = []
    title_bottom = max((r.y1 for r, t in ext_lines if r.y0 < 110), default=40)
    prev_bot = title_bottom + 4
    for k, row in enumerate(rows):
        row_top_est = min(r.y0 for r in row) - 210
        top = max(prev_bot, row_top_est)
        bot = max(r.y1 for r in row) + 62
        bounds.append((top, min(bot, ph)))
        prev_bot = bot + 2
    # fix overlaps with next row's estimate by midpoint
    for k in range(len(bounds) - 1):
        nxt_est = min(r.y0 for r in rows[k + 1]) - 210
        if nxt_est > bounds[k][1]:
            mid = (bounds[k][1] + nxt_est) / 2
            bounds[k] = (bounds[k][0], mid)
            bounds[k + 1] = (mid, bounds[k + 1][1])
    def best_cut_x(top, bot, mid, window=45):
        """Column with least content near mid — avoids slicing logos."""
        band = occ[int(top):int(bot)]
        lo = max(0, int(mid - window)); hi = min(pw, int(mid + window))
        counts = band[:, lo:hi].sum(axis=0)
        m = counts.min()
        cand = np.flatnonzero(counts == m)
        return lo + int(cand[np.abs(cand - (int(mid) - lo)).argmin()])

    blocks = []
    used = np.zeros_like(occ)
    for ridx, ((top, bot), row) in enumerate(zip(bounds, rows)):
        row.sort(key=lambda r: r.x0)
        centers = [(r.x0 + r.x1) / 2 for r in row]
        edges = [0.0]
        for a, b in zip(centers[:-1], centers[1:]):
            edges.append(best_cut_x(top, bot, (a + b) / 2))
        edges.append(float(pw))
        for i in range(len(row)):
            cell = fitz.Rect(edges[i], top, edges[i + 1], bot)
            t = trim_rect(occ, cell)
            if t:
                blocks.append({"rect": t, "banner": False,
                               "key": (top, 1, t.x0)})
                used[int(t.y0):int(t.y1), int(t.x0):int(t.x1)] = True
    # leftover strong content (title banner, section banners)
    leftover = occ & ~used
    ph_, pw_ = leftover.shape
    row_cnt = leftover.sum(axis=1)
    for a, b in runs_of(row_cnt > 3):
        if b - a < 8:
            continue
        band = leftover[a:b]
        density = band.sum() / band.size
        cols_any = band.any(axis=0)
        x0 = int(np.argmax(cols_any)); x1 = int(pw_ - np.argmax(cols_any[::-1]))
        if density < 0.02 or (x1 - x0) < 30:
            continue  # faint decoration
        has_text = any(r.y0 >= a - 4 and r.y1 <= b + 4 for r, t in ext_lines)
        if not has_text:
            continue
        blocks.append({"rect": fitz.Rect(max(0, x0 - PAD), max(0, a - PAD),
                                         min(pw_, x1 + PAD), min(ph_, b + PAD)),
                       "banner": True, "key": (float(a), 0, 0.0)})
    blocks.sort(key=lambda bl: bl["key"])
    return blocks

def build_blocks(page, occ, opts):
    if opts.get("people_grid"):
        pg = people_grid_blocks(page, occ)
        if pg:
            return pg
    ph, pw = occ.shape
    row_any = occ.any(axis=1)
    if not row_any.any():
        return [{"rect": fitz.Rect(page.rect), "banner": False}]
    top, bot = int(np.argmax(row_any)), int(ph - np.argmax(row_any[::-1]))
    cuts = find_cuts(page, occ, opts.get("y_gap", Y_GAP), opts.get("forced", []),
                     opts.get("cut_tol", 2))
    edges = [top] + [c for c in cuts if top < c < bot] + [bot]
    hdil = dilate(occ, 2, axis=1)
    blocks = []
    for b0, b1 in zip(edges[:-1], edges[1:]):
        band = occ[b0:b1]
        if not band.any():
            continue
        rows = band.any(axis=1)
        y0 = b0 + int(np.argmax(rows))
        y1 = b1 - int(np.argmax(rows[::-1]))
        cols_any = hdil[y0:y1].any(axis=0)
        x0, x1 = int(np.argmax(cols_any)), int(pw - np.argmax(cols_any[::-1]))
        cells = None
        h, w = y1 - y0, x1 - x0
        if not opts.get("no_xsplit") and h >= 90 and w / max(h, 1) >= 2.2:
            col_runs = [(a, b) for a, b in runs_of(cols_any[x0:x1])]
            gaps_ok = []
            merged = []
            for a, b in col_runs:
                if merged and a - merged[-1][1] < X_GAP:
                    merged[-1] = (merged[-1][0], b)
                else:
                    merged.append((a, b))
            if len(merged) >= 3:
                widths = [b - a for a, b in merged]
                mean = sum(widths) / len(widths)
                sd = (sum((w_ - mean) ** 2 for w_ in widths) / len(widths)) ** 0.5
                if sd / mean <= 0.6:
                    cells = [(y0, y1, x0 + a, x0 + b) for a, b in merged]
        for (cy0, cy1, cx0, cx1) in (cells or [(y0, y1, x0, x1)]):
            r = fitz.Rect(max(0, cx0 - PAD), max(0, cy0 - PAD),
                          min(pw, cx1 + PAD), min(ph, cy1 + PAD))
            if r.width > 2 and r.height > 2:
                blocks.append({"rect": r, "banner": False})
    return blocks

def pack(blocks):
    """Choose column count & uniform content scale; banner blocks get their
    own scale (fit to column). Returns (scales list, column assignment)."""
    if not blocks:
        return [], []
    base = [b for b in blocks if not b["banner"]] or blocks
    maxw = max(b["rect"].width for b in base)
    maxh = max(b["rect"].height for b in base)
    sumh = sum(b["rect"].height for b in base)
    best = None
    for C in range(1, 7):
        colw = (SW - 2 * MARGIN - (C - 1) * GUTTER) / C
        avail_h = SH - 2 * MARGIN
        s = min(colw / maxw, (C * avail_h * 0.94) / sumh, avail_h / maxh, S_CAP)
        if C > len(blocks):
            s = 0
        if best is None or s > best[1]:
            best = (C, s)
    C, s = best
    colw = (SW - 2 * MARGIN - (C - 1) * GUTTER) / C
    scales = [min(colw / b["rect"].width, s * 1.6, S_CAP) if b["banner"] else s
              for b in blocks]
    target = sum(b["rect"].height * sc for b, sc in zip(blocks, scales)) / C
    cols, acc, cur = [], 0.0, []
    for i, b in enumerate(blocks):
        h = b["rect"].height * scales[i]
        remaining_cols = C - len(cols)
        if cur and acc + h / 2 > target and remaining_cols > 1 and len(blocks) - i >= remaining_cols - 1:
            cols.append(cur); cur, acc = [], 0.0
        cur.append(i); acc += h
    if cur:
        cols.append(cur)
    while len(cols) < C and any(len(c) > 1 for c in cols):
        hi = max(range(len(cols)),
                 key=lambda k: sum(blocks[i]["rect"].height * scales[i] for i in cols[k]))
        col = cols[hi]
        if len(col) < 2:
            break
        half = len(col) // 2
        cols[hi:hi + 1] = [col[:half], col[half:]]
    return scales, cols

def render_page_no_text(doc_path, pno, zoom):
    doc = fitz.open(doc_path)
    page = doc[pno]
    for b in page.get_text("dict")["blocks"]:
        if b["type"] != 0:
            continue
        for l in b["lines"]:
            for s in l["spans"]:
                r = fitz.Rect(s["bbox"]) & page.rect
                if not r.is_empty:
                    page.add_redact_annot(r)
    page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE,
                          graphics=fitz.PDF_REDACT_LINE_ART_NONE,
                          text=fitz.PDF_REDACT_TEXT_REMOVE)
    pm = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    img = Image.frombytes("RGB", (pm.width, pm.height), pm.samples)
    doc.close()
    return img

def dominant_color(img):
    small = np.asarray(img.resize((80, 80), Image.BILINEAR), dtype=np.int32)
    q = (small // 24) * 24 + 12
    colors, counts = np.unique(q.reshape(-1, 3), axis=0, return_counts=True)
    c = colors[counts.argmax()]
    return RGBColor(*(int(min(255, v)) for v in c))

def add_segment(slide, seg, sx, sy, s):
    """Add one text segment; (sx, sy) maps PDF coords -> slide points."""
    bb = seg["bbox"]
    vertical = abs(seg["dir"][1]) > 0.7
    x0, y0 = sx(bb.x0), sy(bb.y0)
    w, h = bb.width * s, bb.height * s
    if vertical:
        cx, cy = x0 + w / 2, y0 + h / 2
        w, h = h, w
        x0, y0 = cx - w / 2, cy - h / 2
    pad_w = w * 0.30 + 4
    tb = slide.shapes.add_textbox(Emu(int(x0 * EMU)), Emu(int(y0 * EMU)),
                                  Emu(int((w + pad_w) * EMU)), Emu(int(max(h, 2) * EMU)))
    if vertical:
        tb.rotation = 270 if seg["dir"][1] < 0 else 90
        tb.left = Emu(int((x0 - pad_w / 2) * EMU))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if any("؀" <= ch <= "ۿ" for r in seg["runs"] for ch in r["text"]):
        p._p.get_or_add_pPr().set("rtl", "1")
        p.alignment = PP_ALIGN.RIGHT
    for r_ in seg["runs"]:
        run = p.add_run()
        run.text = r_["text"].rstrip("\n")
        fam, bold, italic = map_font(r_["font"], r_["flags"])
        run.font.size = Pt(max(1.0, round(r_["size"] * s * 2) / 2))
        run.font.name = fam
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = int_to_rgb(r_["color"])
        if r_["flags"] & 1:
            run.font._rPr.set("baseline", "30000")

def main():
    doc = fitz.open(SRC)
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU))
    prs.slide_height = Emu(int(SH * EMU))
    blank = prs.slide_layouts[6]
    slide_no = 0

    for pno in range(len(doc)):
        if pno in SKIP_PAGES:
            continue
        slide_no += 1
        page = doc[pno]
        zoom = RENDER_LONG_PX / max(page.rect.width, page.rect.height)
        img = render_page_no_text(SRC, pno, zoom)
        slide = prs.slides.add_slide(blank)
        bgf = slide.background.fill
        bgf.solid()
        bgf.fore_color.rgb = dominant_color(img)

        opts = PAGE_OPTS.get(pno, {})
        occ = occupancy(page, img, zoom, opts.get("dist", DIST_THRESH),
                        opts.get("edge_zero", 0))
        blocks = build_blocks(page, occ, opts)
        scales, cols = pack(blocks)
        segments = collect_segments(page)

        # geometry per column
        C = max(1, len(cols))
        colw = (SW - 2 * MARGIN - (C - 1) * GUTTER) / C
        placed = {}  # block index -> (dst_x, dst_y)
        for ci, col in enumerate(cols):
            cx = MARGIN + ci * (colw + GUTTER)
            tot = sum(blocks[i]["rect"].height * scales[i] for i in col)
            free = max(0.0, SH - 2 * MARGIN - tot)
            gap = free / (len(col) + 1)
            y = MARGIN + gap
            for i in col:
                r = blocks[i]["rect"]
                placed[i] = (cx + (colw - r.width * scales[i]) / 2.0, y)
                y += r.height * scales[i] + gap

        for i, bl in enumerate(blocks):
            r = bl["rect"]
            si = scales[i]
            dx, dy = placed[i]
            crop = img.crop((int(r.x0 * zoom), int(r.y0 * zoom),
                             int(r.x1 * zoom), int(r.y1 * zoom)))
            jp, pn = io.BytesIO(), io.BytesIO()
            crop.save(jp, "JPEG", quality=92, optimize=True)
            crop.save(pn, "PNG", optimize=True)
            buf, ext = (pn, "png") if pn.tell() < jp.tell() else (jp, "jpg")
            ipath = os.path.join(WORK, f"blk_{pno+1:02d}_{i:02d}.{ext}")
            with open(ipath, "wb") as f:
                f.write(buf.getvalue())
            slide.shapes.add_picture(ipath, Emu(int(dx * EMU)), Emu(int(dy * EMU)),
                                     Emu(int(r.width * si * EMU)), Emu(int(r.height * si * EMU)))

        # place text of each segment relative to its block
        for seg in segments:
            c = (seg["bbox"].tl + seg["bbox"].br) / 2
            owner = None
            for i, bl in enumerate(blocks):
                if bl["rect"].contains(c):
                    owner = i; break
            if owner is None:
                best_d, owner = 1e18, 0
                for i, bl in enumerate(blocks):
                    bc = (bl["rect"].tl + bl["rect"].br) / 2
                    d = (bc.x - c.x) ** 2 + (bc.y - c.y) ** 2
                    if d < best_d:
                        best_d, owner = d, i
            r = blocks[owner]["rect"]
            si = scales[owner]
            dx, dy = placed[owner]
            sx = lambda x, r=r, dx=dx, si=si: dx + (x - r.x0) * si
            sy = lambda y, r=r, dy=dy, si=si: dy + (y - r.y0) * si
            add_segment(slide, seg, sx, sy, si)

        lo, hi = (min(scales), max(scales)) if scales else (1.0, 1.0)
        print(f"page {pno+1}: {len(blocks)} blocks, {C} cols, scale {lo:.2f}-{hi:.2f}")

    prs.save(OUT)
    print("saved", OUT, f"{os.path.getsize(OUT)/1e6:.1f} MB")

if __name__ == "__main__":
    main()
